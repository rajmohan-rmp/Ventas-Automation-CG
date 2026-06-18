#!/usr/bin/env python3
"""
ventas_slide_generator.py  -  Ventas Daily Dashboard Slide (v3)
Dark-themed dashboard. Accepts multiple HTML files to build a 7-day trend.

Usage (single):
    python ventas_slide_generator.py --html-file body.html --output "Ventas_2026-06-11.pptx"

Usage (7-day trend — pass multiple files, oldest first):
    python ventas_slide_generator.py \
        --html-file day1.html day2.html day3.html ... day7.html \
        --output "Ventas_2026-06-11.pptx"
"""

import argparse, sys, re, tempfile
from pathlib import Path
from datetime import datetime, timedelta

def _ensure(pkg, import_as=None):
    try: __import__(import_as or pkg)
    except ImportError:
        import subprocess
        subprocess.check_call([sys.executable,"-m","pip","install",pkg,
                               "--break-system-packages","-q"])

_ensure("bs4","bs4"); _ensure("matplotlib"); _ensure("python-pptx","pptx")

from bs4 import BeautifulSoup
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Assets ────────────────────────────────────────────────────────────────────
_ROOTS = [
    Path(r"C:\Users\raj.mohan\AppData\Roaming\Claude\local-agent-mode-sessions\skills-plugin\bc21a8f6-04e4-4a84-b221-c011b25b0b70\2471927d-bbdd-47e3-9b37-cca4b89e5be6\skills\6d-design-language\assets"),
    Path("/sessions/funny-happy-mccarthy/mnt/.claude/skills/6d-design-language/assets"),
    Path(__file__).parent / "assets",
]
def _asset(n):
    for r in _ROOTS:
        p = r/n
        try:
            if p.exists(): return p
        except OSError:
            continue
    return _ROOTS[0]/n

# ── Colours ───────────────────────────────────────────────────────────────────
BG         = RGBColor(0x06,0x22,0x4A)
PANEL      = RGBColor(0x0A,0x2A,0x57)
CARD_BLUE  = RGBColor(0x0E,0x51,0x8C)
CARD_TEAL  = RGBColor(0x11,0x72,0xB8)
CARD_RED   = RGBColor(0xC0,0x39,0x2B)
CARD_GREY  = RGBColor(0x34,0x49,0x5E)
CARD_GR    = RGBColor(0x16,0x7E,0xC0)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
LGREY      = RGBColor(0xB0,0xC4,0xDE)
CHART_BG   = "#0A2A57"
CHART_BG2  = "#0E3A6E"
B1,B2,B3,B4 = "#1172B8","#17A2B8","#5BC0DE","#27AE60"
LYELLOW     = "#F1C40F"   # total line — yellow
COLORS_POOL = [B1,B2,B3,B4,"#E67E22"]
# Completed=light green, InProgress=light orange, Pending=red; extra slots for ST
STACK_COLORS= ["#2ECC71","#F39C12","#E74C3C","#9B59B6","#3498DB"]

# ── Parsing ───────────────────────────────────────────────────────────────────
def parse_html(html):
    soup = BeautifulSoup(html,"html.parser")
    result = {}
    for h3,tbl in zip(soup.find_all("h3"),soup.find_all("table")):
        rows=[]
        for tr in tbl.find_all("tr"):
            cells=[td.get_text(strip=True) for td in tr.find_all(["th","td"])]
            if cells: rows.append(cells)
        result[h3.get_text(strip=True)]=rows
    m=re.search(r"Report Date\s*:\s*(\d{4}-\d{2}-\d{2})",html)
    result["report_date"]=m.group(1) if m else datetime.today().strftime("%Y-%m-%d")
    return result

def merge_data(html_list):
    """Merge tables from multiple daily emails; keep last 7 unique dates per table."""
    combined={}
    rdates=[]
    for html in html_list:
        d=parse_html(html)
        rd=d.pop("report_date",None)
        if rd: rdates.append(rd)
        for key,rows in d.items():
            if not rows: continue
            if key not in combined:
                combined[key]={"headers":rows[0],"rows":{}}
            for r in rows[1:]:
                if r: combined[key]["rows"][r[0]]=r  # keyed by date — later overrides earlier
    # Re-assemble: sort by date, keep last 7
    result={}
    for key,v in combined.items():
        sorted_rows=sorted(v["rows"].values(),key=lambda r:r[0])
        result[key]=[v["headers"]]+sorted_rows[-7:]
    result["report_date"]=max(rdates) if rdates else datetime.today().strftime("%Y-%m-%d")
    return result

def _val(rows,date,col):
    if not rows: return 0
    h=rows[0]
    try: idx=h.index(col)
    except ValueError: return 0
    vals=[int(r[idx]) for r in rows[1:]
          if len(r)>idx and (date is None or r[0]==date) and r[idx].isdigit()]
    return sum(vals)

def _dates(rows):
    if not rows or len(rows)<2: return []
    seen=[];out=[]
    for r in rows[1:]:
        if r[0] not in seen: seen.append(r[0]);out.append(r[0])
    return out

def _short(d):
    try: return datetime.strptime(d,"%Y-%m-%d").strftime("%d-%b")
    except: return d

# ── CSV ingestion (latest mail's attachments) ─────────────────────────────────
import csv as _csv

def _norm(k):
    return (k or "").strip().upper().replace(" ", "_").lstrip("﻿")

def _read_csv_rows(path):
    raw=open(path,encoding="utf-8-sig",errors="ignore").read()
    if not raw.strip(): return []
    # delimiter sniff
    first=raw.splitlines()[0]
    delim=","
    for d in [",",";","\t","|"]:
        if d in first: delim=d; break
    rdr=_csv.DictReader(raw.splitlines(),delimiter=delim)
    out=[]
    for r in rdr:
        out.append({_norm(k):(v.strip() if isinstance(v,str) else v) for k,v in r.items()})
    return out

def _to_int(v):
    try: return int(float(str(v).strip()))
    except: return 0

def _date_key(row):
    for k in ("ORDER_DATE","PAYMENT_DATE","DATE","REPORT_DATE","TXN_DATE"):
        if k in row and row[k]: return row[k][:10]
    # fallback: first column that looks like a date
    for k,v in row.items():
        if v and re.match(r"\d{4}-\d{2}-\d{2}", str(v)): return str(v)[:10]
    return None

def _status_bucket(s, mapping):
    su=_norm(s)
    for frag,target in mapping:
        if frag in su: return target
    return None

def _agg_sales(rows):
    """Return {date: [completed,inprogress,pending,total]} summed across order types."""
    agg={}
    has_summary = rows and ("COMPLETED" in rows[0])
    has_status  = rows and ("STATUS" in rows[0])
    smap=[("COMPLET","COMPLETED"),("PROGRESS","INPROGRESS"),("PEND","PENDING"),
          ("SUCCESS","COMPLETED"),("FAIL","PENDING")]
    for r in rows:
        d=_date_key(r)
        if not d: continue
        a=agg.setdefault(d,[0,0,0,0])
        if has_summary:
            c=_to_int(r.get("COMPLETED")); i=_to_int(r.get("INPROGRESS")); pp=_to_int(r.get("PENDING"))
            t=_to_int(r.get("TOTAL_COUNT")) or (c+i+pp)
            a[0]+=c; a[1]+=i; a[2]+=pp; a[3]+=t
        elif has_status:
            b=_status_bucket(r.get("STATUS"),smap)
            if b=="COMPLETED": a[0]+=1
            elif b=="INPROGRESS": a[1]+=1
            elif b=="PENDING": a[2]+=1
            a[3]+=1
        else:
            a[3]+=1
    return agg

def _agg_stock(rows):
    """Return {date: [approved,delivered,gr_pending,created,total]}."""
    agg={}
    has_summary = rows and ("APPROVED" in rows[0] or "DELIVERED" in rows[0])
    has_status  = rows and ("STATUS" in rows[0])
    smap=[("APPROV","APPROVED"),("DELIVER","DELIVERED"),("GR","GR_PENDING"),
          ("GOODS","GR_PENDING"),("CREAT","CREATED"),("NEW","CREATED")]
    for r in rows:
        d=_date_key(r)
        if not d: continue
        a=agg.setdefault(d,[0,0,0,0,0])
        if has_summary:
            ap=_to_int(r.get("APPROVED")); de=_to_int(r.get("DELIVERED"))
            gr=_to_int(r.get("GR_PENDING")); cr=_to_int(r.get("CREATED"))
            t=_to_int(r.get("TOTAL_REQUESTS")) or (ap+de+gr+cr)
            a[0]+=ap; a[1]+=de; a[2]+=gr; a[3]+=cr; a[4]+=t
        elif has_status:
            b=_status_bucket(r.get("STATUS"),smap)
            idx={"APPROVED":0,"DELIVERED":1,"GR_PENDING":2,"CREATED":3}.get(b)
            if idx is not None: a[idx]+=1
            a[4]+=1
        else:
            a[4]+=1
    return agg

def _classify(path):
    n=Path(path).name.upper()
    if "PRIMARY" in n: return "Primary Sales Summary"
    if "POS"     in n: return "POS Sales Summary"
    if "STOCK" in n or "TRANSFER" in n: return "Stock Transfer Summary"
    return None

def csv_to_data(csv_paths, report_date=""):
    """Build the same data dict the renderer expects, from the latest mail's CSVs.
    Aggregates per ORDER_DATE across all order types (fixes undercount)."""
    data={}
    all_dates=[]
    for path in csv_paths:
        key=_classify(path)
        if not key:
            print(f"WARN: could not classify {path} by filename; skipping")
            continue
        rows=_read_csv_rows(path)
        if not rows:
            print(f"WARN: {path} is empty"); continue
        print(f"  {Path(path).name}: {len(rows)} rows, cols={list(rows[0].keys())}")
        if key=="Stock Transfer Summary":
            agg=_agg_stock(rows)
            hdr=["ORDER_DATE","APPROVED","DELIVERED","GR_PENDING","CREATED","TOTAL_REQUESTS"]
        else:
            agg=_agg_sales(rows)
            hdr=["ORDER_DATE","COMPLETED","INPROGRESS","PENDING","TOTAL_COUNT"]
        out=[hdr]
        for d in sorted(agg):
            out.append([d]+[str(x) for x in agg[d]])
            all_dates.append(d)
        data[key]=out
    rd=report_date or (max(all_dates) if all_dates else datetime.today().strftime("%Y-%m-%d"))
    # Keep only the fixed 7-day window ending on report_date so KPI 7-DAY TOTAL
    # matches the chart window.
    try:
        end=datetime.strptime(rd,"%Y-%m-%d")
        win={(end-timedelta(days=i)).strftime("%Y-%m-%d") for i in range(1,8)}  # exclude report day
        for key in list(data.keys()):
            if key=="report_date": continue
            tbl=data[key]
            data[key]=[tbl[0]]+[r for r in tbl[1:] if r[0] in win]
    except Exception:
        pass
    data["report_date"]=rd
    return data

# ── Chart ─────────────────────────────────────────────────────────────────────
def make_chart(title, rows, col_bars, col_line, tmp_dir, report_date=None, palette=None):
    # Always show a 7-day window; zero-fill dates with no data
    if report_date:
        try:
            end=datetime.strptime(report_date,"%Y-%m-%d")-timedelta(days=1)  # exclude report day
            dates=[(end-timedelta(days=6-i)).strftime("%Y-%m-%d") for i in range(7)]
        except:
            dates=_dates(rows) or ["No data"]
    else:
        dates=_dates(rows) or ["No data"]
    x=np.arange(len(dates))
    bar_data={c:[_val(rows,d,c) for d in dates] for c in col_bars}
    line_data=[_val(rows,d,col_line) for d in dates]

    fig,ax=plt.subplots(figsize=(7.5,6.5))  # matches reduced panel aspect ratio
    fig.patch.set_facecolor(CHART_BG)
    ax.set_facecolor(CHART_BG2)

    # Pre-compute ylim so bar label threshold is available during plotting
    ymax_line=max(line_data) if any(v>0 for v in line_data) else 1
    ylim_top=max(ymax_line*1.2,1)

    n=len(col_bars)
    colors=(palette or STACK_COLORS)[:n]
    bottom=np.zeros(len(dates))
    bar_w=0.55
    min_seg=max((ylim_top/0.9)*0.05,0.5)  # min segment height to show label
    for (col,vals),color in zip(bar_data.items(),colors):
        vals_arr=np.array(vals,dtype=float)
        ax.bar(x,vals_arr,bar_w,bottom=bottom,label=col,color=color,zorder=3,alpha=0.93)
        txt_col='#1A1A1A' if color in ("#F1C40F","#F39C12","#FFD700","#2ECC71","#FFC107") else 'white'
        for i,v in enumerate(vals):
            if v>0:
                if v>=min_seg:
                    ax.text(x[i],bottom[i]+v/2,str(v),
                            ha='center',va='center',color=txt_col,fontsize=10,fontweight='bold')
                else:
                    # bar too thin — show label above with a dark pill so it's always readable
                    ax.text(x[i],bottom[i]+v+ylim_top*0.02,str(v),
                            ha='center',va='bottom',color='white',fontsize=9,fontweight='bold',
                            bbox=dict(boxstyle='round,pad=0.18',facecolor='#1a1a2e',
                                      edgecolor=color,linewidth=0.8,alpha=0.9))
        bottom+=vals_arr

    ax2=ax.twinx()
    ax2.plot(x,line_data,color=LYELLOW,linewidth=2.0,marker='o',markersize=5,zorder=4)
    for xi,v in zip(x,line_data):
        if v>0:
            ax2.annotate(str(v),(xi,v),textcoords="offset points",xytext=(0,6),
                         ha='center',color='white',fontsize=11,fontweight='bold')

    ax.set_xticks(x)
    ax.set_xticklabels([_short(d) for d in dates],color='white',fontsize=11,rotation=20,ha='right')
    ax.tick_params(colors='white'); ax.set_yticks([]); ax2.set_yticks([])
    ax2.tick_params(colors='white'); ax.spines[:].set_visible(False); ax2.spines[:].set_visible(False)
    ax.set_title(title,color='white',fontsize=13,fontweight='bold',pad=8)
    ax.grid(axis='y',color='white',alpha=0.06,zorder=0)
    ax.set_ylim(0, ylim_top/0.9)   # bar scale 10% higher → bars appear 10% below line
    ax2.set_ylim(0, ylim_top)

    handles=[plt.Rectangle((0,0),1,1,color=c) for c in colors]
    handles+=[plt.Line2D([0],[0],color=LYELLOW,lw=2,marker='o',ms=4)]
    ax.legend(handles,list(col_bars)+["Total"],
              loc='upper center',bbox_to_anchor=(0.5,-0.06),
              facecolor=CHART_BG,edgecolor='none',labelcolor='white',
              fontsize=8,ncol=n+1,framealpha=0.5,
              handlelength=1.0,handleheight=0.8,borderpad=0.4,columnspacing=0.8)

    plt.tight_layout(pad=0.4)
    plt.subplots_adjust(bottom=0.16)
    path=tmp_dir/f"{title.replace(' ','_')}.png"
    fig.savefig(str(path),dpi=200,bbox_inches='tight',facecolor=CHART_BG)
    plt.close(fig)
    return path

# ── Full-dashboard PNG (matplotlib only — no LibreOffice needed) ───────────────
def _draw_chart_ax(ax, title, rows, col_bars, col_line, report_date=None, palette=None):
    if report_date:
        try:
            end=datetime.strptime(report_date,"%Y-%m-%d")-timedelta(days=1)
            dates=[(end-timedelta(days=6-i)).strftime("%Y-%m-%d") for i in range(7)]
        except Exception:
            dates=_dates(rows) or ["No data"]
    else:
        dates=_dates(rows) or ["No data"]
    x=np.arange(len(dates))
    bar_data={c:[_val(rows,d,c) for d in dates] for c in col_bars}
    line_data=[_val(rows,d,col_line) for d in dates]
    ax.set_facecolor(CHART_BG2)
    ymax_line=max(line_data) if any(v>0 for v in line_data) else 1
    ylim_top=max(ymax_line*1.2,1)
    n=len(col_bars); colors=(palette or STACK_COLORS)[:n]
    bottom=np.zeros(len(dates)); bar_w=0.55
    min_seg=max((ylim_top/0.9)*0.05,0.5)
    for (col,vals),color in zip(bar_data.items(),colors):
        vals_arr=np.array(vals,dtype=float)
        ax.bar(x,vals_arr,bar_w,bottom=bottom,label=col,color=color,zorder=3,alpha=0.93)
        txt_col='#1A1A1A' if color in ("#F1C40F","#F39C12","#FFD700","#2ECC71","#FFC107") else 'white'
        for i,vv in enumerate(vals):
            if vv>0:
                if vv>=min_seg:
                    ax.text(x[i],bottom[i]+vv/2,str(vv),ha='center',va='center',color=txt_col,fontsize=8,fontweight='bold')
                else:
                    ax.text(x[i],bottom[i]+vv+ylim_top*0.02,str(vv),ha='center',va='bottom',color='white',fontsize=7,fontweight='bold',
                            bbox=dict(boxstyle='round,pad=0.18',facecolor='#1a1a2e',edgecolor=color,linewidth=0.8,alpha=0.9))
        bottom+=vals_arr
    ax2=ax.twinx()
    ax2.plot(x,line_data,color=LYELLOW,linewidth=2.0,marker='o',markersize=4,zorder=4)
    for xi,vv in zip(x,line_data):
        if vv>0:
            ax2.annotate(str(vv),(xi,vv),textcoords="offset points",xytext=(0,5),ha='center',color='white',fontsize=8,fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels([_short(d) for d in dates],color='white',fontsize=8,rotation=20,ha='right')
    ax.tick_params(colors='white'); ax.set_yticks([]); ax2.set_yticks([])
    ax2.tick_params(colors='white'); ax.spines[:].set_visible(False); ax2.spines[:].set_visible(False)
    ax.set_title(title,color='white',fontsize=11,fontweight='bold',pad=6)
    ax.grid(axis='y',color='white',alpha=0.06,zorder=0)
    ax.set_ylim(0,ylim_top/0.9); ax2.set_ylim(0,ylim_top)
    handles=[plt.Rectangle((0,0),1,1,color=c) for c in colors]
    handles+=[plt.Line2D([0],[0],color=LYELLOW,lw=2,marker='o',ms=4)]
    ax.legend(handles,list(col_bars)+["Total"],loc='upper center',bbox_to_anchor=(0.5,-0.08),
              facecolor=CHART_BG,edgecolor='none',labelcolor='white',fontsize=6.5,ncol=n+1,framealpha=0.5,
              handlelength=1.0,handleheight=0.8,borderpad=0.3,columnspacing=0.6)

def render_dashboard_png(data, report_date, out_png, dpi=150):
    """Render the whole dashboard (header + KPI tiles + 3 charts) as one PNG,
    using matplotlib only. No LibreOffice/PowerPoint required."""
    from matplotlib.patches import Rectangle
    rd=report_date or data.get("report_date")
    st =data.get("Stock Transfer Summary",[]); pos=data.get("POS Sales Summary",[]); pri=data.get("Primary Sales Summary",[])
    try: last_day=(datetime.strptime(rd,"%Y-%m-%d")-timedelta(days=1)).strftime("%Y-%m-%d")
    except Exception: last_day=rd
    v=_val; TEAL="#1172B8"; RED="#C0392B"; GREY="#34495E"
    def trip(rows,d): return [("Completed" if False else "", )]
    pri_td=[("Completed",v(pri,None,"COMPLETED"),TEAL),("InProgress",v(pri,None,"INPROGRESS"),GREY),("Pending",v(pri,None,"PENDING"),RED)]
    pri_ld=[("Completed",v(pri,last_day,"COMPLETED"),TEAL),("InProgress",v(pri,last_day,"INPROGRESS"),GREY),("Pending",v(pri,last_day,"PENDING"),RED)]
    pos_td=[("Completed",v(pos,None,"COMPLETED"),TEAL),("InProgress",v(pos,None,"INPROGRESS"),GREY),("Pending",v(pos,None,"PENDING"),RED)]
    pos_ld=[("Completed",v(pos,last_day,"COMPLETED"),TEAL),("InProgress",v(pos,last_day,"INPROGRESS"),GREY),("Pending",v(pos,last_day,"PENDING"),RED)]
    st_td=[("Orders",v(st,None,"TOTAL_REQUESTS"),GREY),("Delivered",v(st,None,"DELIVERED"),TEAL),("GR Pending",v(st,None,"GR_PENDING"),RED)]
    st_ld=[("Orders",v(st,last_day,"TOTAL_REQUESTS"),GREY),("Delivered",v(st,last_day,"DELIVERED"),TEAL),("GR Pending",v(st,last_day,"GR_PENDING"),RED)]
    sections=[("Primary Sales",pri_td,pri_ld),("POS Sales",pos_td,pos_ld),("Stock Transfer",st_td,st_ld)]
    chart_specs=[("Primary Sales (7-Day Trend)",pri,["COMPLETED","INPROGRESS","PENDING"],"TOTAL_COUNT",None),
                 ("POS Sales (7-Day Trend)",pos,["COMPLETED","INPROGRESS","PENDING"],"TOTAL_COUNT",None),
                 ("Stock Transfer (7-Day Trend)",st,["APPROVED","DELIVERED","GR_PENDING","CREATED"],"TOTAL_REQUESTS",
                  ["#FFC107","#2ECC71","#E74C3C","#9B59B6","#3498DB"])]
    fig=plt.figure(figsize=(13.333,7.5),dpi=dpi); fig.patch.set_facecolor("#06224A")
    bg=fig.add_axes([0,0,1,1]); bg.set_axis_off(); bg.set_xlim(0,1); bg.set_ylim(0,1)
    bg.add_patch(Rectangle((0,0.92),1,0.08,color="#0E518C",zorder=1))
    bg.text(0.015,0.965,"VENTAS DAILY DASHBOARD",ha="left",va="center",color="white",fontsize=20,fontweight="bold",zorder=3)
    bg.text(0.015,0.933,"Report Date: %s    |    Airtel Congo (CG)    |    Last 7 Days Trend"%_short_date_long(rd),
            ha="left",va="center",color="#B0C4DE",fontsize=9.5,zorder=3)
    bg.text(0.985,0.96,"6D Technologies",ha="right",va="center",color="white",fontsize=13,fontweight="bold",zorder=3)
    M=0.015; colw=(1-4*M)/3; ptop=0.90
    g=0.006; sh_h=0.030; rowlabel_h=0.018; tile_h=0.072; rg=0.008
    def row_tiles(x,y,w,h,tiles):
        tg=0.006; tw=(w-2*tg)/3
        for j,(lbl,val,color) in enumerate(tiles):
            tx=x+j*(tw+tg)
            bg.add_patch(Rectangle((tx,y),tw,h,color=color,zorder=2))
            bg.text(tx+tw/2,y+h*0.66,lbl,ha="center",va="center",color="white",fontsize=7.5,fontweight="bold",zorder=3)
            bg.text(tx+tw/2,y+h*0.30,str(val),ha="center",va="center",color="white",fontsize=14,fontweight="bold",zorder=3)
    for i,(name,td,ld) in enumerate(sections):
        x0=M+i*(colw+M)
        panel_h=g+sh_h+g+rowlabel_h+tile_h+rg+rowlabel_h+tile_h+g
        pbot=ptop-panel_h
        bg.add_patch(Rectangle((x0,pbot),colw,panel_h,color="#0A2A57",zorder=1))
        bg.add_patch(Rectangle((x0+g,ptop-g-sh_h),colw-2*g,sh_h,color="#0E518C",zorder=2))
        bg.text(x0+colw/2,ptop-g-sh_h/2,"%s  (as of %s)"%(name,_short(last_day)),
                ha="center",va="center",color="white",fontsize=8.5,fontweight="bold",zorder=3)
        y=ptop-g-sh_h-g
        bg.text(x0+g,y-rowlabel_h/2,"7-DAY TOTAL",ha="left",va="center",color="white",fontsize=7,fontweight="bold",zorder=3)
        y=y-rowlabel_h; row_tiles(x0+g,y-tile_h,colw-2*g,tile_h,td); y=y-tile_h-rg
        bg.text(x0+g,y-rowlabel_h/2,"LAST DAY",ha="left",va="center",color="white",fontsize=7,fontweight="bold",zorder=3)
        y=y-rowlabel_h; row_tiles(x0+g,y-tile_h,colw-2*g,tile_h,ld)
    cb_bottom=0.10; cb_h=0.35
    for i,(title,rows,bars,line,palette) in enumerate(chart_specs):
        x0=M+i*(colw+M)+0.012
        ax=fig.add_axes([x0,cb_bottom,colw-0.024,cb_h])
        _draw_chart_ax(ax,title,rows,bars,line,rd,palette)
    fig.savefig(str(out_png),dpi=dpi,facecolor=fig.get_facecolor())
    plt.close(fig)
    print("Dashboard PNG rendered (matplotlib):",out_png)
    return out_png

# ── PPTX helpers ─────────────────────────────────────────────────────────────
def _bg(slide,color):
    f=slide.background.fill; f.solid(); f.fore_color.rgb=color

def _rect(slide,x,y,w,h,fc=None,lc=None,lw=0,rx=0):
    st=MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
    s=slide.shapes.add_shape(st,int(x),int(y),int(w),int(h))
    if rx: s.adjustments[0]=rx
    if fc: s.fill.solid(); s.fill.fore_color.rgb=fc
    else:  s.fill.background()
    if lc: s.line.color.rgb=lc; s.line.width=Pt(lw or 1)
    else:  s.line.fill.background()
    return s

def _txt(slide,text,x,y,w,h,sz=11,bold=False,col=WHITE,align=PP_ALIGN.LEFT,italic=False):
    tb=slide.shapes.add_textbox(int(x),int(y),int(w),int(h))
    tf=tb.text_frame; tf.word_wrap=False
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=str(text)
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col; r.font.italic=italic
    return tb

def kpi_tile(slide,label,value,x,y,w,h,tc,vsz=18):
    _rect(slide,x,y,w,h,fc=tc,rx=0.05)
    # Single textbox spanning full tile — XML vertical-center anchors both lines
    PAD=Inches(0.04)
    tb=slide.shapes.add_textbox(int(x+PAD),int(y),int(w-2*PAD),int(h))
    tf=tb.text_frame; tf.word_wrap=False
    ns='http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr=tf._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is not None: bodyPr.set('anchor','ctr')
    # Paragraph 1 — label
    p1=tf.paragraphs[0]; p1.alignment=PP_ALIGN.CENTER
    r1=p1.add_run(); r1.text=label
    r1.font.size=Pt(8.5); r1.font.bold=True; r1.font.color.rgb=WHITE
    # Paragraph 2 — value
    p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
    r2=p2.add_run(); r2.text=str(value)
    r2.font.size=Pt(vsz); r2.font.bold=True; r2.font.color.rgb=WHITE


def sec_hdr(slide,label,x,y,w):
    _rect(slide,x,y,w,Inches(0.22),fc=CARD_BLUE,rx=0.04)
    # Full-height textbox + XML vertical-center so text sits in middle of strip
    tb=slide.shapes.add_textbox(int(x+Inches(0.05)),int(y),int(w-Inches(0.10)),int(Inches(0.22)))
    tf=tb.text_frame; tf.word_wrap=False
    ns='http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr=tf._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is not None: bodyPr.set('anchor','ctr')
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=label
    r.font.size=Pt(9); r.font.bold=True; r.font.color.rgb=WHITE

def _row_lbl(slide,text,x,y,w,h):
    tb=slide.shapes.add_textbox(int(x),int(y),int(w),int(h))
    tf=tb.text_frame; tf.word_wrap=False
    ns='http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr=tf._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is not None: bodyPr.set('anchor','ctr')
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=text
    r.font.size=Pt(8.5); r.font.bold=True; r.font.color.rgb=WHITE

# ── Main ──────────────────────────────────────────────────────────────────────
def generate_slide(html_list, report_date, output_path, data=None, png_path=None):
    if data is None:
        if isinstance(html_list, str): html_list=[html_list]
        data=merge_data(html_list)
    if report_date: data["report_date"]=report_date
    rd=data["report_date"]

    st =data.get("Stock Transfer Summary",[])
    pos=data.get("POS Sales Summary",[])
    pri=data.get("Primary Sales Summary",[])

    # KPI "LAST DAY" = report day - 1 (the report-generation day is excluded)
    try: last_day=(datetime.strptime(rd,"%Y-%m-%d")-timedelta(days=1)).strftime("%Y-%m-%d")
    except: last_day=rd
    st_td=pos_td=pri_td=last_day

    # Last-day KPIs
    st_ord=_val(st,st_td,"TOTAL_REQUESTS"); st_del=_val(st,st_td,"DELIVERED"); st_gr=_val(st,st_td,"GR_PENDING")
    pos_ok=_val(pos,pos_td,"COMPLETED"); pos_pend=_val(pos,pos_td,"PENDING"); pos_inprog=_val(pos,pos_td,"INPROGRESS")
    pri_ok=_val(pri,pri_td,"COMPLETED"); pri_pend=_val(pri,pri_td,"PENDING"); pri_inprog=_val(pri,pri_td,"INPROGRESS")
    # 7-day totals (date=None → sum all rows in window)
    pri_ok7=_val(pri,None,"COMPLETED"); pri_ip7=_val(pri,None,"INPROGRESS"); pri_pend7=_val(pri,None,"PENDING")
    pos_ok7=_val(pos,None,"COMPLETED"); pos_ip7=_val(pos,None,"INPROGRESS"); pos_pend7=_val(pos,None,"PENDING")
    st_ord7=_val(st,None,"TOTAL_REQUESTS"); st_del7=_val(st,None,"DELIVERED"); st_gr7=_val(st,None,"GR_PENDING")

    with tempfile.TemporaryDirectory() as td_dir:
        tmp=Path(td_dir)
        pri_chart=make_chart("Primary Sales (7-Day Trend)",
                             pri,["COMPLETED","INPROGRESS","PENDING"],"TOTAL_COUNT",tmp,rd)
        pos_chart=make_chart("POS Sales (7-Day Trend)",
                             pos,["COMPLETED","INPROGRESS","PENDING"],"TOTAL_COUNT",tmp,rd)
        st_chart =make_chart("Stock Transfer (7-Day Trend)",
                             st,["APPROVED","DELIVERED","GR_PENDING","CREATED"],"TOTAL_REQUESTS",tmp,rd,
                             palette=["#FFC107","#2ECC71","#E74C3C","#9B59B6","#3498DB"])

        prs=Presentation()
        prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
        SW=prs.slide_width; SH=prs.slide_height
        slide=prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide,BG)

        # ── Header ──────────────────────────────────────────────────────────
        HDR_H=Inches(0.55)
        bp=_asset("banner-gradient.png")
        if bp.exists(): slide.shapes.add_picture(str(bp),0,0,SW,HDR_H)
        else: _rect(slide,0,0,SW,HDR_H,fc=CARD_BLUE)

        _txt(slide,"VENTAS DAILY DASHBOARD",
             Inches(0.25),Inches(0.05),Inches(9),Inches(0.34),sz=20,bold=True)
        _txt(slide,f"Report Date: {_short_date_long(rd)}   |   Airtel Congo (CG)   |   Last 7 Days Trend",
             Inches(0.25),Inches(0.38),Inches(9),Inches(0.20),sz=8.5,col=LGREY)
        lg=_asset("6d-logo-white.png")
        if lg.exists():
            p=slide.shapes.add_picture(str(lg),0,Inches(0.07),height=Inches(0.38))
            p.left=SW-p.width-Inches(0.2)

        # ── KPI row  (Last Day + 7-Day Total) ───────────────────────────────
        KY=Inches(0.62); M=Inches(0.17)
        G=Inches(0.07); SW3=(SW-4*M)/3; TW3=(SW3-4*G)/3
        HDRS=Inches(0.22)   # section header strip
        SLH =Inches(0.18)   # sub-row label height ("7-DAY TOTAL" / "LAST DAY")
        TH  =Inches(0.52)   # tile height — must be ≥ (label+value) to contain text
        RG  =Inches(0.05)   # gap between the two tile rows
        KH  =G+HDRS+G+SLH+TH+RG+SLH+TH+G  # ≈1.64"

        def sx(i): return M+i*(SW3+M)

        def _kpi_row(x, ld_vals, td_vals, lbls, cols):
            """Draw 7-Day-Total row (top) then Last-Day row (bottom)."""
            for row_vals, row_tag, row_y in [
                (td_vals, "7-DAY TOTAL",  KY+G+HDRS+G),
                (ld_vals, "LAST DAY",     KY+G+HDRS+G+SLH+TH+RG),
            ]:
                _row_lbl(slide, row_tag, x+G, row_y, SW3-2*G, SLH)
                ty = row_y + SLH
                for j,(lbl,val,col) in enumerate(zip(lbls,row_vals,cols)):
                    kpi_tile(slide,lbl,val, x+G+j*(TW3+G), ty, TW3, TH, col, vsz=20)

        # Primary Sales
        x0=sx(0)
        _rect(slide,x0,KY,SW3,KH,fc=PANEL,rx=0.05)
        sec_hdr(slide,f"Primary Sales  (as of {_short_date_long(pri_td)})",x0+G,KY+G,SW3-2*G)
        _kpi_row(x0,
                 [pri_ok, pri_inprog, pri_pend],
                 [pri_ok7,pri_ip7,    pri_pend7],
                 ["Completed","InProgress","Pending"],
                 [CARD_TEAL,CARD_GREY,CARD_RED])

        # POS Sales
        x1=sx(1)
        _rect(slide,x1,KY,SW3,KH,fc=PANEL,rx=0.05)
        sec_hdr(slide,f"POS Sales  (as of {_short_date_long(pos_td)})",x1+G,KY+G,SW3-2*G)
        _kpi_row(x1,
                 [pos_ok, pos_inprog, pos_pend],
                 [pos_ok7,pos_ip7,    pos_pend7],
                 ["Completed","InProgress","Pending"],
                 [CARD_TEAL,CARD_GREY,CARD_RED])

        # Stock Transfer
        x2=sx(2)
        _rect(slide,x2,KY,SW3,KH,fc=PANEL,rx=0.05)
        sec_hdr(slide,f"Stock Transfer  (as of {_short_date_long(st_td)})",x2+G,KY+G,SW3-2*G)
        _kpi_row(x2,
                 [st_ord, st_del, st_gr],
                 [st_ord7,st_del7,st_gr7],
                 ["Orders","Delivered","GR Pending"],
                 [CARD_GREY,CARD_TEAL,CARD_RED])

        # ── Charts (3 panels) ─────────────────────────────────────────────────
        CY=KY+KH+Inches(0.10); CH=(SH-CY-Inches(0.20))*0.80; CW=(SW-4*M)/3
        _rect(slide,M,            CY,CW,CH,fc=PANEL,rx=0.04)
        _rect(slide,M+CW+M,       CY,CW,CH,fc=PANEL,rx=0.04)
        _rect(slide,M+2*(CW+M),   CY,CW,CH,fc=PANEL,rx=0.04)

        PAD=Inches(0.12)
        if pri_chart.exists():
            slide.shapes.add_picture(str(pri_chart),
                int(M+PAD),int(CY+PAD),width=int(CW-2*PAD))
        if pos_chart.exists():
            slide.shapes.add_picture(str(pos_chart),
                int(M+CW+M+PAD),int(CY+PAD),width=int(CW-2*PAD))
        if st_chart.exists():
            slide.shapes.add_picture(str(st_chart),
                int(M+2*(CW+M)+PAD),int(CY+PAD),width=int(CW-2*PAD))

        # ── Footer ───────────────────────────────────────────────────────────
        _rect(slide,0,SH-Inches(0.27),SW,Inches(0.27),fc=CARD_BLUE)
        _txt(slide,
             f"Generated by Ventas USDM 2.0  |  {datetime.now().strftime('%d %b %Y, %H:%M')} IST",
             Inches(0.3),SH-Inches(0.24),Inches(9),Inches(0.20),sz=7.5,col=LGREY)
        _txt(slide,"CONFIDENTIAL",SW-Inches(2.4),SH-Inches(0.24),Inches(2.3),Inches(0.20),
             sz=7.5,col=LGREY,align=PP_ALIGN.RIGHT,italic=True)

        prs.save(output_path)
        print(f"Saved: {output_path}")
        if png_path:
            try:
                render_dashboard_png(data, rd, png_path)
            except Exception as e:
                print("WARN: dashboard PNG render failed:", e)

    return output_path

def _short_date_long(d):
    try: return datetime.strptime(d,"%Y-%m-%d").strftime("%d-%b-%Y")
    except: return d

if __name__=="__main__":
    parser=argparse.ArgumentParser(description="Ventas daily dashboard PPTX (7-day trend)")
    parser.add_argument("--html-file",nargs="+",help="One or more HTML files (oldest to newest)")
    parser.add_argument("--csv-file",nargs="+",help="CSV attachments from the LATEST mail (PRIMARY/POS/STOCK)")
    parser.add_argument("--html",help="Single HTML string")
    parser.add_argument("--date",default="",help="Report date YYYY-MM-DD")
    parser.add_argument("--output",required=True,help="Output .pptx path")
    parser.add_argument("--png",default="",help="Also render the dashboard as a PNG at this path")
    args=parser.parse_args()
    if args.csv_file:
        print("Building from latest mail CSV attachments:")
        data=csv_to_data(args.csv_file, args.date)
        generate_slide(None, args.date, args.output, data=data, png_path=(args.png or None))
    elif args.html_file:
        html_list=[Path(f).read_text(encoding="utf-8") for f in args.html_file]
        generate_slide(html_list, args.date, args.output, png_path=(args.png or None))
    elif args.html:
        generate_slide([args.html], args.date, args.output, png_path=(args.png or None))
    else:
        parser.error("Provide --csv-file, --html-file, or --html")
